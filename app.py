from flask import Flask, request, render_template, send_file, redirect, url_for, flash
import os
import pptx
from pptx.util import Inches, Pt
import google.generativeai as genai
from imagescrapper import get_image

genai.configure(api_key='AIzaSyBwOXEqZPXHy7UXa3oQDP9gIGDIHoVaMdk')

model = genai.GenerativeModel('gemini-pro')

#font sizes for title and content
TITLE_FONT_SIZE = Pt(30)
SLIDE_FONT_SIZE = Pt(16)

app = Flask(__name__)

#upload and dowload paths
upload_directory = 'uploads'
output_directory = os.path.join(os.path.expanduser("~"), "Downloads")

#Create directories if not exist
os.makedirs(upload_directory, exist_ok=True)
os.makedirs(output_directory, exist_ok=True)

@app.route('/', methods=['GET', 'POST'])
def index():
    ppt_filename = None  #ppt_filename variable

    if request.method == 'POST':
        topic = request.form.get('topic')
        num_slides = request.form.get("num_slides")

        if topic:
            slide_titles = generate_slide_titles(topic, num_slides)
            slide_contents = [generate_slide_content(title) for title in slide_titles]
            image_subjects = [generate_image_subject(content) for content in slide_contents]
            print(image_subjects)

            image_paths = ['' if subject == 'None' else get_image(subject, 1)[0] for subject in image_subjects]
            print(image_paths)

            #Generating presentation file
            ppt_filename = os.path.join(output_directory, f"{topic}_presentation.pptx")
            create_presentation(topic, slide_titles, slide_contents, image_paths, ppt_filename)

    # Render the HTML form and pass the ppt_filename variable to the template
    return render_template('index.html', ppt_filename=ppt_filename)


@app.route('/download/<filename>')
def download(filename):
    #download path
    filepath = os.path.join(output_directory, filename)
    
    #Check if the file exists
    if not os.path.exists(filepath):
        return f"File not found: {filepath}", 404
    #Return the file for download
    return send_file(filepath, as_attachment=True, download_name=filename)


def generate_slide_titles(topic, num_slides):
    # Generate slide titles using the Gemini API
    prompt = f"Generate {num_slides} slide titles that are structured and relevant to the topic: {topic}. Each title should be concise and informative."
    response = model.generate_content(prompt)
    slide_titles = response.text.split('\n')
    slide_titles = [title.strip() for title in slide_titles if title.strip()]#breaking up response into slides
    for i in range(len(slide_titles)):
        slide_titles[i] = slide_titles[i].replace("**", "") #stripping markdown tags
    return slide_titles

def generate_slide_content(slide_title): #for each slide title from the list, generate content
  prompt = f"Generate detailed content for the slide titled: {slide_title}. The content should be informative and structured, suitable for a presentation slide. Content should strictly be within 100 words and have only 4 major points maximum. do not include the slide title again in this response.no markdown"
  response = model.generate_content(prompt)
  return response.text.strip().replace("**", "")

def create_presentation(topic, slide_titles, slide_contents, image_paths, ppt_filename):
    
    prs = pptx.Presentation("themes\\theme1.pptx")#create object, using existing theme

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = topic.upper()
    subtitle.text = "A Short Presentation"

    for i, (slide_title, slide_content) in enumerate(zip(slide_titles, slide_contents)):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_title
        slide.shapes.placeholders[1].text = slide_content

        slide.shapes.title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE

        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = SLIDE_FONT_SIZE
        
        print(i)
        image_path = image_paths[i] if i < len(image_paths) else None
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(6.1), Inches(4.8), width=Inches(3), height=Inches(2))
    #removing the extra slide at beginning    
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    xml_slides.remove(slides[0]) 

    prs.save(ppt_filename)

def generate_image_subject(slide_content):
    #Generate an image subject based on the slide content.
    prompt = f"Generate an image subject based on the following slide content: {slide_content}. The image should aesthetic or apt enough to be in a presentation.it should be small and concise enough for a search engine to understand.return one subject only if image is needed in slide - if not needed, return None"
    response = model.generate_content(prompt)
    return response.text.strip()

if __name__ == '__main__':
    app.run(debug=True)